from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from slides_maker.domain.schemas import Slide


class PPTBuilder:
    def __init__(self) -> None:
        self.presentation = Presentation()

    def add_slide(self, slide: Slide, image_path: Path | None = None) -> None:
        if slide.slide_type == "title":
            self._add_title_slide(slide)
        elif slide.slide_type == "section":
            self._add_section_slide(slide)
        else:
            self._add_content_slide(slide, image_path=image_path)

    def save(self, path: Path) -> None:
        path.parent.mkdir(parents=True, exist_ok=True)
        self.presentation.save(str(path))

    def _add_title_slide(self, slide: Slide) -> None:
        layout = self.presentation.slide_layouts[0]
        s = self.presentation.slides.add_slide(layout)
        s.shapes.title.text = slide.heading
        if len(slide.bullet_points) > 0:
            subtitle = s.placeholders[1]
            subtitle.text = "\n".join(slide.bullet_points)

    def _add_section_slide(self, slide: Slide) -> None:
        layout = self.presentation.slide_layouts[2]
        s = self.presentation.slides.add_slide(layout)
        s.shapes.title.text = slide.heading
        body = s.placeholders[1].text_frame
        body.clear()
        for point in slide.bullet_points:
            p = body.add_paragraph()
            p.text = point
            p.level = 0

    def _add_content_slide(self, slide: Slide, image_path: Path | None) -> None:
        layout = self.presentation.slide_layouts[5]
        s = self.presentation.slides.add_slide(layout)
        title = s.shapes.title
        title.text = slide.heading

        left = Inches(0.7)
        top = Inches(1.6)
        width = Inches(5.4)
        height = Inches(4.6)
        textbox = s.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        for i, point in enumerate(slide.bullet_points):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = point
            p.level = 0
            p.font.size = Pt(20)

        if image_path is not None and image_path.exists():
            s.shapes.add_picture(str(image_path), Inches(6.4), Inches(1.6), width=Inches(3.1))


def build_presentation(slides: Iterable[Slide], image_paths: list[Path]) -> Presentation:
    builder = PPTBuilder()
    for slide, img in zip(slides, image_paths):
        builder.add_slide(slide, image_path=img)
    return builder.presentation
